"""
Corpus ingestion: from a folder of real documents to training-ready text.

An organization's knowledge does not arrive as prompt/response pairs. It
arrives as thousands of PDFs, Word files, web exports, and images. This module
is the front of the pipeline that turns that pile into something a stratum can
train on:

  ingest()          walk a folder, extract text from every supported file
                    (images via a vision teacher), deduplicate, optionally
                    redact, and cut the text into overlapping chunks with
                    full provenance - which file, which position, which hash
  generate_pairs()  have a teacher model write question/answer pairs grounded
                    in each chunk, split into train and held-out test sets

Built for corpora in the thousands of files:

  - every extraction is CACHED by the file's content hash, so re-running an
    ingest touches only new or changed files - a crash costs nothing
  - a corrupt or unreadable file is recorded and skipped, never fatal
  - pair generation appends as it goes, retries flaky teachers, and resumes
    mid-corpus exactly like teacher-gen (docs/07)
  - the train/test split is decided per chunk by a stable hash, so no chunk
    ever leaks between train and test across resumed or repeated runs

What this deliberately does NOT do: turn the corpus into the model's memory.
Fine-tuning teaches skills and style reliably, facts unreliably - the
knowledge layer for "answer questions about our documents" is retrieval
(RAG), explained in docs/14-from-corpus-to-model.md. This module builds the
SKILL data: extraction, terminology, formats, grounded Q&A behavior.
"""
from __future__ import annotations

import hashlib
import json
import re
import time
from pathlib import Path

DOC_TYPES = {".txt", ".md", ".rst", ".csv", ".log", ".html", ".htm", ".pdf",
             ".docx", ".pptx", ".xlsx"}
IMAGE_TYPES = {".png", ".jpg", ".jpeg", ".webp", ".gif", ".bmp", ".tiff"}


class CorpusError(Exception):
    """A per-file extraction problem, with a message a non-ML developer can act on."""


_CTYPE_EXT = {"text/html": ".html", "application/pdf": ".pdf",
              "text/plain": ".txt", "text/markdown": ".md"}


def _filename_for(url: str, content_type: str) -> str:
    from urllib.parse import unquote, urlsplit
    name = unquote(urlsplit(url).path.rstrip("/").rsplit("/", 1)[-1]) or "download"
    name = re.sub(r"[^\w.() -]", "_", name)[:120]
    if "." not in name:
        name += _CTYPE_EXT.get(content_type, ".txt")
    return name


def fetch_urls(urls: list[str], out_dir: str, retries: int = 3,
               verbose: bool = True) -> dict:
    """Download a list of URLs into a folder ready for ingest.

    The acquisition step of the pipeline: give it the pages, reports, and
    files the corpus should contain, and it pulls them down with the same
    survival rules as everything else here - retries with growing pauses,
    a failed URL recorded and skipped rather than fatal, and re-running
    skips files already downloaded. Filenames come from the URL, with the
    extension corrected from the server's content type so ingest can
    dispatch the right extractor.
    """
    import urllib.request

    from . import __version__

    out = Path(out_dir)
    out.mkdir(parents=True, exist_ok=True)
    counts = {"fetched": 0, "already": 0, "failed": 0}
    for url in urls:
        url = url.strip()
        if not url or url.startswith("#"):
            continue
        # Resume check before any network traffic: the name is known up front
        # except for its extension, so probe the possible extensions too.
        stem = _filename_for(url, "")
        candidates = {stem} | {stem.rsplit(".", 1)[0] + ext
                               for ext in _CTYPE_EXT.values()}
        existing = next((c for c in candidates if (out / c).exists()), None)
        if existing:
            counts["already"] += 1
            if verbose:
                print(f" already have {existing}")
            continue

        data = None
        error = None
        ctype = ""
        for attempt in range(retries):
            try:
                req = urllib.request.Request(
                    url, headers={"User-Agent": f"stratum-corpus/{__version__}"})
                with urllib.request.urlopen(req, timeout=60) as resp:
                    ctype = (resp.headers.get("Content-Type") or "").split(";")[0].strip()
                    data = resp.read()
                break
            except Exception as e:
                error = e
                wait = 2 ** attempt
                if verbose:
                    print(f" {url}: {e} - retrying in {wait}s ({attempt + 1}/{retries})")
                time.sleep(wait)
        if data is None:
            counts["failed"] += 1
            if verbose:
                print(f" FAILED {url}: {error}")
            continue
        name = _filename_for(url, ctype)
        (out / name).write_bytes(data)
        counts["fetched"] += 1
        if verbose:
            print(f" fetched {name} ({len(data):,} bytes)")
    if verbose:
        print(f"\nFetched {counts['fetched']} files "
              f"({counts['already']} already present, {counts['failed']} failed) "
              f"-> {out_dir}")
        if counts["failed"]:
            print(" failed URLs can simply be re-run - existing files are kept.")
    return counts


def _extract_txt(path: Path) -> str:
    return path.read_text(encoding="utf-8", errors="replace")


def _extract_html(path: Path) -> str:
    from html.parser import HTMLParser

    class TextGrabber(HTMLParser):
        def __init__(self):
            super().__init__()
            self.parts = []
            self.skip = 0

        def handle_starttag(self, tag, attrs):
            if tag in ("script", "style"):
                self.skip += 1

        def handle_endtag(self, tag):
            if tag in ("script", "style") and self.skip:
                self.skip -= 1
            if tag in ("p", "div", "br", "li", "tr", "h1", "h2", "h3", "h4"):
                self.parts.append("\n")

        def handle_data(self, data):
            if not self.skip:
                self.parts.append(data)

    grabber = TextGrabber()
    grabber.feed(path.read_text(encoding="utf-8", errors="replace"))
    return "".join(grabber.parts)


def _extract_pdf(path: Path) -> str:
    try:
        from pypdf import PdfReader
    except ImportError:
        raise CorpusError(
            "PDF support needs pypdf. Install the corpus extras: "
            "pip install 'stratum-slm[corpus]'"
        )
    reader = PdfReader(str(path))
    pages = []
    for i, page in enumerate(reader.pages):
        text = page.extract_text() or ""
        if text.strip():
            pages.append(text)
    text = "\n\n".join(pages)
    if not text.strip():
        raise CorpusError(
            "no extractable text - this is probably a scanned PDF (pictures of "
            "pages). Export its pages as images and ingest those with a vision "
            "teacher (--images), or run OCR before ingesting."
        )
    return text


def _extract_docx(path: Path) -> str:
    try:
        import docx
    except ImportError:
        raise CorpusError(
            "Word support needs python-docx. Install the corpus extras: "
            "pip install 'stratum-slm[corpus]'"
        )
    document = docx.Document(str(path))
    parts = [p.text for p in document.paragraphs]
    for table in document.tables:
        for row in table.rows:
            parts.append(" | ".join(cell.text for cell in row.cells))
    return "\n".join(parts)


def _extract_pptx(path: Path) -> str:
    try:
        from pptx import Presentation
    except ImportError:
        raise CorpusError(
            "PowerPoint support needs python-pptx. Install the corpus extras: "
            "pip install 'stratum-slm[corpus]'"
        )
    deck = Presentation(str(path))
    parts = []
    for i, slide in enumerate(deck.slides, 1):
        parts.append(f"\nSlide {i}:")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs).strip()
                    if line:
                        parts.append(line)
            if shape.has_table:
                for row in shape.table.rows:
                    parts.append(" | ".join(cell.text for cell in row.cells))
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            note = slide.notes_slide.notes_text_frame.text.strip()
            if note:
                parts.append(f"Speaker notes: {note}")
    return "\n".join(parts)


def _extract_xlsx(path: Path) -> str:
    try:
        import openpyxl
    except ImportError:
        raise CorpusError(
            "Excel support needs openpyxl. Install the corpus extras: "
            "pip install 'stratum-slm[corpus]'"
        )
    book = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
    parts = []
    for sheet in book.worksheets:
        parts.append(f"\nSheet {sheet.title}:")
        for row in sheet.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None]
            if cells:
                parts.append(" | ".join(cells))
    book.close()
    return "\n".join(parts)


def extract_text(path: Path) -> str:
    """Extract plain text from one document. Raises CorpusError with a fix."""
    ext = path.suffix.lower()
    if ext in (".html", ".htm"):
        return _extract_html(path)
    if ext == ".pdf":
        return _extract_pdf(path)
    if ext == ".docx":
        return _extract_docx(path)
    if ext == ".pptx":
        return _extract_pptx(path)
    if ext == ".xlsx":
        return _extract_xlsx(path)
    return _extract_txt(path)


_REDACTIONS = [
    ("email", re.compile(r"[\w.+-]+@[\w-]+\.[\w.-]+")),
    ("phone", re.compile(r"(?<![\d-])(?:\+?\d[\d\s().-]{8,}\d)(?![\d-])")),
    ("card", re.compile(r"\b(?:\d[ -]?){13,16}\b")),
]


def redact(text: str) -> tuple[str, dict]:
    """Replace obvious emails, phone numbers, and card-like numbers with tags.

    A baseline, not a guarantee - names, addresses, and domain-specific
    identifiers need the organization's own rules. Real compliance work runs
    its own scrubber BEFORE ingest and treats this as a second net.
    """
    counts = {}
    for name, pattern in _REDACTIONS:
        text, n = pattern.subn(f"[{name} removed]", text)
        if n:
            counts[name] = n
    return text, counts


def chunk_text(text: str, size: int = 2400, overlap: int = 240) -> list[tuple[int, str]]:
    """Cut text into overlapping chunks of about `size` characters.

    A chunk is the unit a teacher can actually read and write pairs about -
    a whole 300-page manual is too much, a sentence too little. Overlap keeps
    facts that straddle a boundary intact in at least one chunk. Cuts prefer
    a paragraph break, then a sentence end, near the target size.
    """
    if size <= overlap:
        raise ValueError("chunk size must be larger than the overlap")
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text).strip()
    chunks = []
    start = 0
    while start < len(text):
        end = min(start + size, len(text))
        if end < len(text):
            window = text[start:end]
            cut = max(window.rfind("\n\n"), window.rfind(". "))
            if cut > size // 2:
                end = start + cut + 1
        piece = text[start:end].strip()
        if piece:
            chunks.append((start, piece))
        if end >= len(text):
            break
        start = end - overlap
    return chunks


def _sha256_file(path: Path) -> str:
    h = hashlib.sha256()
    with open(path, "rb") as f:
        for block in iter(lambda: f.read(1 << 20), b""):
            h.update(block)
    return h.hexdigest()


def _sha256_text(text: str) -> str:
    return hashlib.sha256(text.encode("utf-8")).hexdigest()


def ingest(in_dir: str, out_dir: str, vision_teacher=None, redact_pii: bool = False,
           chunk_size: int = 2400, overlap: int = 240, verbose: bool = True) -> dict:
    """Walk a corpus folder and produce chunks.jsonl plus a manifest.

    Every file's extracted text is cached under out_dir/cache keyed by the
    file's content hash - re-runs only extract what is new or changed, which
    is what makes a many-thousand-file corpus practical (and what makes the
    expensive vision extractions one-time costs). chunks.jsonl is rebuilt
    fresh each run from the caches, so deletions and edits are always
    reflected. Duplicate files (same content, any name) are chunked once.

    vision_teacher: a callable image-path -> text, from stratum.vision. When
    None, images are counted and skipped with a note.

    Returns a summary dict (also the last line printed).
    """
    inp = Path(in_dir)
    if not inp.is_dir():
        raise FileNotFoundError(f"Corpus folder not found: {in_dir}")
    out = Path(out_dir)
    cache = out / "cache"
    cache.mkdir(parents=True, exist_ok=True)

    manifest = []
    seen_hashes = set()
    chunks_path = out / "chunks.jsonl"
    n_chunks = 0
    counts = {"documents": 0, "images": 0, "cached": 0, "duplicates": 0,
              "skipped": 0, "errors": 0}

    files = sorted(p for p in inp.rglob("*")
                   if p.is_file() and not any(part.startswith(".") for part in p.parts))

    with open(chunks_path, "w", encoding="utf-8") as chunks_file:
        for path in files:
            rel = str(path.relative_to(inp))
            ext = path.suffix.lower()
            entry = {"path": rel, "kind": None, "status": "ok", "chunks": 0}

            if ext in DOC_TYPES:
                entry["kind"] = "document"
            elif ext in IMAGE_TYPES:
                entry["kind"] = "image"
                if vision_teacher is None:
                    entry["status"] = "skipped"
                    entry["reason"] = ("image skipped - pass --images to extract its "
                                      "content with a vision teacher")
                    counts["skipped"] += 1
                    manifest.append(entry)
                    continue
            else:
                entry["kind"] = "other"
                entry["status"] = "skipped"
                entry["reason"] = f"unsupported type '{ext}'"
                counts["skipped"] += 1
                manifest.append(entry)
                continue

            sha = _sha256_file(path)
            entry["sha256"] = sha
            if sha in seen_hashes:
                entry["status"] = "duplicate"
                counts["duplicates"] += 1
                manifest.append(entry)
                continue
            seen_hashes.add(sha)

            cached = cache / f"{sha}.txt"
            try:
                if cached.exists():
                    text = cached.read_text(encoding="utf-8")
                    counts["cached"] += 1
                else:
                    if entry["kind"] == "image":
                        if verbose:
                            print(f" reading image {rel} with the vision teacher")
                        text = vision_teacher(str(path))
                    else:
                        text = extract_text(path)
                    cached.write_text(text, encoding="utf-8")
                counts["documents" if entry["kind"] == "document" else "images"] += 1
            except CorpusError as e:
                entry["status"] = "error"
                entry["reason"] = str(e)
                counts["errors"] += 1
                manifest.append(entry)
                if verbose:
                    print(f" ERROR {rel}: {e}")
                continue
            except Exception as e:
                entry["status"] = "error"
                entry["reason"] = f"{type(e).__name__}: {e}"
                counts["errors"] += 1
                manifest.append(entry)
                if verbose:
                    print(f" ERROR {rel}: {e}")
                continue

            if redact_pii:
                text, removed = redact(text)
                if removed:
                    entry["redacted"] = removed

            pieces = chunk_text(text, chunk_size, overlap)
            entry["chunks"] = len(pieces)
            for start, piece in pieces:
                chunks_file.write(json.dumps({
                    "id": _sha256_text(piece)[:16],
                    "text": piece,
                    "source": rel,
                    "source_sha256": sha,
                    "start_char": start,
                    "kind": entry["kind"],
                }, ensure_ascii=False) + "\n")
                n_chunks += 1
            manifest.append(entry)

    (out / "manifest.jsonl").write_text(
        "\n".join(json.dumps(e, ensure_ascii=False) for e in manifest) + "\n",
        encoding="utf-8")

    summary = dict(counts)
    summary["chunks"] = n_chunks
    summary["files"] = len(files)
    if verbose:
        print(f"\nIngested {counts['documents']} documents and {counts['images']} images "
              f"({counts['cached']} from cache) -> {n_chunks} chunks")
        print(f" duplicates: {counts['duplicates']}  skipped: {counts['skipped']}  "
              f"errors: {counts['errors']}")
        print(f" chunks -> {chunks_path}")
        print(f" per-file record -> {out / 'manifest.jsonl'}")
        if counts["errors"]:
            print(" errors are listed in the manifest with a reason each - fix or "
                  "remove those files and re-run, everything else is cached.")
    return summary


PAIR_PROMPT = """\
You are writing training data for a smaller model.

Write {n} question-and-answer pairs grounded ONLY in the passage below.
Rules:
- Every question must stand alone: name the subject explicitly, never write
  "this document", "the passage", or "the text above".
- Every answer must be fully supported by the passage - no outside knowledge.
- {instruction}

Return ONLY a JSON array, nothing else, in exactly this shape:
[{{"prompt": "...", "response": "..."}}]

Passage (from {source}):
{chunk}
"""


def _parse_pairs(raw: str) -> list[dict]:
    """Pull prompt/response pairs out of whatever a teacher wrote.

    The prompt asks for one JSON array, but real teachers - small local ones
    especially - also emit one object per line, prose around the JSON, or
    trailing commentary after the array. This scans the text and collects
    every JSON value that carries a prompt and a response, so a chunk only
    fails when the output contains no usable pairs at all.
    """
    decoder = json.JSONDecoder()
    good = []
    i = 0
    while i < len(raw):
        if raw[i] not in "[{":
            i += 1
            continue
        try:
            value, end = decoder.raw_decode(raw, i)
        except ValueError:
            i += 1
            continue
        items = value if isinstance(value, list) else [value]
        for p in items:
            if isinstance(p, dict) and p.get("prompt") and p.get("response"):
                good.append({"prompt": str(p["prompt"]),
                             "response": str(p["response"])})
        i = end
    if not good:
        raise ValueError("teacher output contained no usable prompt/response pairs")
    return good


def _is_test_chunk(chunk_id: str, test_fraction: float, seed: int) -> bool:
    # Stable per chunk: the same chunk lands on the same side of the split in
    # every run and every resume, so test content never leaks into training.
    h = hashlib.sha256(f"{seed}:{chunk_id}".encode()).hexdigest()
    return (int(h[:8], 16) / 0xFFFFFFFF) < test_fraction


def generate_pairs(chunks_path: str, instruction: str, teacher_fn,
                   out_train: str, out_test: str | None = None,
                   per_chunk: int = 3, test_fraction: float = 0.1,
                   max_chunks: int | None = None,
                   retries: int = 3, seed: int = 42, verbose: bool = True) -> dict:
    """Ask a teacher to write grounded pairs for every chunk.

    Same survival rules as teacher-gen: pairs are appended the moment they
    exist, failed teacher calls retry with growing pauses, and re-running the
    same command resumes - chunks already answered in the output files are
    skipped. The split is per chunk (not per pair), decided by a stable hash,
    so a chunk's content is never in both train and test.

    max_chunks caps teacher cost on a big corpus by sampling evenly across
    the chunk file - every source contributes, rather than the first N chunks
    all coming from whichever file sorts first. The sample is deterministic,
    so re-runs resume cleanly.
    """
    import math

    from .data import load_jsonl

    chunks = load_jsonl(chunks_path, required_keys=("id", "text", "source"))
    if max_chunks and len(chunks) > max_chunks:
        step = math.ceil(len(chunks) / max_chunks)
        chunks = chunks[::step]
        if verbose:
            print(f"Sampling {len(chunks)} of the chunks (every {step}th) "
                  f"to cap teacher cost - raise --max-chunks for more coverage")
    if out_test is None and test_fraction > 0:
        raise ValueError("test_fraction is set but no out_test path was given")

    done = set()
    for p in (out_train, out_test):
        if p and Path(p).exists():
            for line in Path(p).read_text(encoding="utf-8").splitlines():
                if line.strip():
                    done.add(json.loads(line).get("source_chunk"))
    if done and verbose:
        print(f"Resuming: {len(done)} chunks already answered")

    Path(out_train).parent.mkdir(parents=True, exist_ok=True)
    if out_test:
        Path(out_test).parent.mkdir(parents=True, exist_ok=True)

    counts = {"train_pairs": 0, "test_pairs": 0, "chunks_done": len(done),
              "chunks_failed": 0}
    train_f = open(out_train, "a", encoding="utf-8")
    test_f = open(out_test, "a", encoding="utf-8") if out_test else None
    try:
        for i, chunk in enumerate(chunks, 1):
            if chunk["id"] in done:
                continue
            prompt = PAIR_PROMPT.format(n=per_chunk, instruction=instruction,
                                        source=chunk["source"], chunk=chunk["text"])
            pairs = None
            for attempt in range(retries):
                try:
                    pairs = _parse_pairs(teacher_fn(prompt))
                    break
                except Exception as e:
                    wait = 2 ** attempt
                    if verbose:
                        print(f" chunk {i}/{len(chunks)}: {e} - retrying in {wait}s "
                              f"({attempt + 1}/{retries})")
                    time.sleep(wait)
            if pairs is None:
                counts["chunks_failed"] += 1
                continue

            to_test = test_fraction > 0 and _is_test_chunk(chunk["id"], test_fraction, seed)
            for pair in pairs[:per_chunk]:
                record = {"source_chunk": chunk["id"], "source": chunk["source"]}
                if to_test:
                    record.update({"prompt": pair["prompt"],
                                   "expected": pair["response"]})
                    test_f.write(json.dumps(record, ensure_ascii=False) + "\n")
                    counts["test_pairs"] += 1
                else:
                    record.update(pair)
                    train_f.write(json.dumps(record, ensure_ascii=False) + "\n")
                    counts["train_pairs"] += 1
            (test_f or train_f).flush()
            train_f.flush()
            counts["chunks_done"] += 1
            if verbose and counts["chunks_done"] % 25 == 0:
                print(f" {counts['chunks_done']}/{len(chunks)} chunks done")
    finally:
        train_f.close()
        if test_f:
            test_f.close()

    if verbose:
        print(f"\nWrote {counts['train_pairs']} training pairs -> {out_train}")
        if out_test:
            print(f"Wrote {counts['test_pairs']} held-out test pairs -> {out_test}")
        if counts["chunks_failed"]:
            print(f"WARNING: {counts['chunks_failed']} chunks failed after "
                  f"{retries} tries each. Re-run the same command to retry "
                  f"just those.")
    return counts
