"""Tests for the corpus pipeline: extraction, chunking, redaction, ingest
with caching and dedup, pair generation with a stable train/test split, and
the full walk from a raw corpus folder to a trained model."""
import json
from pathlib import Path

import pytest

from stratum.corpus import (CorpusError, chunk_text, extract_text,
                            generate_pairs, ingest, redact)
from stratum.vision import get_vision_teacher

REPORT = ("Pipeline segment 7 was inspected on 14 March. Wall thickness "
          "measured 11.2 mm against a 12 mm nominal. Corrosion allowance "
          "remains within limits. Next inspection due in 24 months. " * 8)


def _write_minimal_pdf(path, text):
    """A smallest-possible one-page PDF with real extractable text."""
    header = b"%PDF-1.4\n"
    objects = [
        b"1 0 obj\n<< /Type /Catalog /Pages 2 0 R >>\nendobj\n",
        b"2 0 obj\n<< /Type /Pages /Kids [3 0 R] /Count 1 >>\nendobj\n",
        (b"3 0 obj\n<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] "
         b"/Contents 4 0 R /Resources << /Font << /F1 5 0 R >> >> >>\nendobj\n"),
        None,  # content stream, built below
        (b"5 0 obj\n<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>"
         b"\nendobj\n"),
    ]
    stream = f"BT /F1 12 Tf 72 712 Td ({text}) Tj ET".encode()
    objects[3] = (b"4 0 obj\n<< /Length " + str(len(stream)).encode() +
                  b" >>\nstream\n" + stream + b"\nendstream\nendobj\n")

    body = b""
    offsets = []
    for obj in objects:
        offsets.append(len(header) + len(body))
        body += obj
    xref_pos = len(header) + len(body)
    xref = b"xref\n0 6\n0000000000 65535 f \n"
    for off in offsets:
        xref += f"{off:010d} 00000 n \n".encode()
    trailer = (b"trailer\n<< /Size 6 /Root 1 0 R >>\nstartxref\n" +
               str(xref_pos).encode() + b"\n%%EOF\n")
    path.write_bytes(header + body + xref + trailer)


@pytest.fixture()
def corpus_dir(tmp_path):
    """A miniature of the client's folder: mixed formats, a duplicate, an
    image, an unsupported file, and a corrupt one."""
    import docx
    from PIL import Image

    root = tmp_path / "corpus"
    root.mkdir()
    (root / "notes.txt").write_text(REPORT, encoding="utf-8")
    (root / "copy of notes.txt").write_text(REPORT, encoding="utf-8")  # duplicate
    (root / "page.html").write_text(
        "<html><script>var x=1</script><body><h1>Safety bulletin</h1>"
        "<p>Valve V-201 requires torque checks every 6 months.</p>"
        "</body></html>", encoding="utf-8")
    d = docx.Document()
    d.add_paragraph("Compressor station C-4 operates at 60 bar.")
    d.add_paragraph("Contact the duty engineer at duty@example.com "
                    "or +1 555 010 4477.")
    d.save(str(root / "manual.docx"))
    _write_minimal_pdf(root / "spec.pdf",
                       "Design pressure for header H-2 is 84 bar.")

    from pptx import Presentation
    deck = Presentation()
    slide = deck.slides.add_slide(deck.slide_layouts[1])
    slide.shapes.title.text = "Turbine maintenance briefing"
    slide.placeholders[1].text = "Blade inspection interval is 8000 hours."
    deck.save(str(root / "briefing.pptx"))

    import openpyxl
    book = openpyxl.Workbook()
    sheet = book.active
    sheet.title = "Readings"
    sheet.append(["Station", "Pressure_bar"])
    sheet.append(["C-4", 60])
    sheet.append(["H-2", 84])
    book.save(str(root / "readings.xlsx"))
    Image.new("RGB", (32, 32), color=(200, 30, 30)).save(root / "gauge.png")
    (root / "data.xyz").write_text("binary-ish", encoding="utf-8")
    (root / "broken.pdf").write_bytes(b"%PDF-1.4 this is not really a pdf")
    return root


def test_extract_each_format(corpus_dir):
    assert "Pipeline segment 7" in extract_text(corpus_dir / "notes.txt")
    html = extract_text(corpus_dir / "page.html")
    assert "Valve V-201" in html
    assert "var x=1" not in html  # scripts stripped
    assert "Compressor station C-4" in extract_text(corpus_dir / "manual.docx")
    assert "84 bar" in extract_text(corpus_dir / "spec.pdf")
    pptx = extract_text(corpus_dir / "briefing.pptx")
    assert "8000 hours" in pptx and "Slide 1" in pptx
    xlsx = extract_text(corpus_dir / "readings.xlsx")
    assert "C-4 | 60" in xlsx


def test_scanned_pdf_gives_actionable_error(tmp_path):
    from pypdf import PdfWriter
    p = tmp_path / "scan.pdf"
    writer = PdfWriter()
    writer.add_blank_page(width=612, height=792)  # a page with no text layer
    with open(p, "wb") as f:
        writer.write(f)
    with pytest.raises(CorpusError, match="scanned"):
        extract_text(p)


def test_chunking_covers_text_with_overlap():
    text = "word " * 2000
    chunks = chunk_text(text, size=1000, overlap=100)
    assert len(chunks) > 5
    for (s1, c1), (s2, c2) in zip(chunks, chunks[1:]):
        assert s2 < s1 + 1000          # neighbors overlap
        assert len(c1) <= 1000
    joined = "".join(c for _, c in chunks)
    assert joined.count("word") >= 2000  # nothing lost


def test_chunking_rejects_bad_sizes():
    with pytest.raises(ValueError):
        chunk_text("abc", size=100, overlap=100)


def test_redact_baseline():
    text, counts = redact("Mail duty@example.com or call +1 555 010 4477.")
    assert "duty@example.com" not in text
    assert "555" not in text
    assert counts == {"email": 1, "phone": 1}


def test_ingest_full_folder(corpus_dir, tmp_path):
    out = tmp_path / "out"
    summary = ingest(str(corpus_dir), str(out),
                     vision_teacher=get_vision_teacher("echo"),
                     redact_pii=True, verbose=False)
    assert summary["documents"] == 6       # txt, html, docx, pdf, pptx, xlsx
    assert summary["images"] == 1
    assert summary["duplicates"] == 1      # the copied txt
    assert summary["skipped"] == 1         # the .xyz file
    assert summary["errors"] == 1          # the broken pdf
    assert summary["chunks"] >= 5

    chunks = [json.loads(l) for l in
              (out / "chunks.jsonl").read_text(encoding="utf-8").splitlines()]
    # Provenance on every chunk.
    for c in chunks:
        assert c["source"] and c["id"] and "start_char" in c
    # The image's extracted content made it in.
    assert any("gauge.png" in c["text"] for c in chunks)
    # Redaction reached the docx content.
    assert not any("duty@example.com" in c["text"] for c in chunks)

    manifest = [json.loads(l) for l in
                (out / "manifest.jsonl").read_text(encoding="utf-8").splitlines()]
    broken = next(m for m in manifest if m["path"] == "broken.pdf")
    assert broken["status"] == "error" and broken["reason"]


def test_ingest_resume_uses_cache(corpus_dir, tmp_path):
    out = tmp_path / "out"
    ingest(str(corpus_dir), str(out), vision_teacher=get_vision_teacher("echo"),
           verbose=False)
    calls = []

    def counting_teacher(path):
        calls.append(path)
        return "should never be called"

    again = ingest(str(corpus_dir), str(out), vision_teacher=counting_teacher,
                   verbose=False)
    assert calls == []                     # image came from cache
    assert again["cached"] == again["documents"] + again["images"]


def _fake_pair_teacher(prompt):
    # A well-behaved teacher: three grounded pairs as a JSON array.
    return json.dumps([
        {"prompt": f"Question {i} about the passage", "response": f"Answer {i}"}
        for i in range(3)])


def test_generate_pairs_split_and_resume(corpus_dir, tmp_path):
    out = tmp_path / "out"
    ingest(str(corpus_dir), str(out), vision_teacher=get_vision_teacher("echo"),
           verbose=False)
    train = tmp_path / "train.jsonl"
    test = tmp_path / "test.jsonl"

    first = generate_pairs(str(out / "chunks.jsonl"), "Ask factual questions.",
                           _fake_pair_teacher, str(train), str(test),
                           test_fraction=0.34, verbose=False)
    assert first["train_pairs"] + first["test_pairs"] > 0
    assert first["chunks_failed"] == 0

    # Chunk-level split: no chunk id appears on both sides.
    train_ids = {json.loads(l)["source_chunk"]
                 for l in train.read_text(encoding="utf-8").splitlines()}
    test_ids = {json.loads(l)["source_chunk"]
                for l in test.read_text(encoding="utf-8").splitlines()}
    assert not (train_ids & test_ids)

    # Test rows carry expected, train rows carry response.
    if test_ids:
        row = json.loads(test.read_text(encoding="utf-8").splitlines()[0])
        assert "expected" in row and "response" not in row

    # Resuming asks the teacher nothing new.
    calls = []

    def counting(prompt):
        calls.append(prompt)
        return _fake_pair_teacher(prompt)

    generate_pairs(str(out / "chunks.jsonl"), "Ask factual questions.",
                   counting, str(train), str(test),
                   test_fraction=0.34, verbose=False)
    assert calls == []


def test_generate_pairs_survives_bad_teacher(corpus_dir, tmp_path, monkeypatch):
    import time as time_mod
    monkeypatch.setattr(time_mod, "sleep", lambda s: None)
    out = tmp_path / "out"
    ingest(str(corpus_dir), str(out), verbose=False)
    counts = generate_pairs(str(out / "chunks.jsonl"), "Ask.",
                            lambda p: "I refuse to emit JSON",
                            str(tmp_path / "t.jsonl"), None,
                            test_fraction=0.0, retries=2, verbose=False)
    assert counts["chunks_failed"] > 0
    assert counts["train_pairs"] == 0


def test_vision_backends():
    teacher = get_vision_teacher("echo")
    assert "gauge.png" in teacher("some/dir/gauge.png")
    with pytest.raises(ValueError, match="bogus"):
        get_vision_teacher("bogus")


def test_corpus_to_trained_model(corpus_dir, tmp_path, tiny_base):
    """The whole front half, end to end: raw folder -> chunks -> pairs ->
    a trained stratum on the tiny model. No gaps between the stages."""
    from stratum.train import train_tile

    out = tmp_path / "out"
    ingest(str(corpus_dir), str(out), vision_teacher=get_vision_teacher("echo"),
           verbose=False)
    train = tmp_path / "train.jsonl"
    generate_pairs(str(out / "chunks.jsonl"), "Ask factual questions.",
                   _fake_pair_teacher, str(train), None,
                   test_fraction=0.0, verbose=False)

    loss = train_tile(skill_path=str(train), out_dir=str(tmp_path / "stratum"),
                      base_model=tiny_base, rank=2, epochs=1, batch_size=2,
                      max_len=96, load_4bit=False, seed=5)
    assert loss == loss
    assert (tmp_path / "stratum" / "stratum_card.json").exists()
